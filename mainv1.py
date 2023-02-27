from datetime import datetime
from itertools import dropwhile, takewhile
import calendar
import instaloader
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import pptx_replace
import os

def findDay(date):
  year, month, day= (int(i) for i in date.split('-'))   

  dayNumber = calendar.weekday(year, month, day)
  days =["Senin", "Selasa", "Rabu", "Kamis","Jumat", "Sabtu", "Minggu"]

  return (days[dayNumber])
    
def findMonth(date):
  year, month, day= (int(i) for i in date.split('-'))

  monthNumber = calendar.month_name(month)
  months = ["Januari","Februari","Maret","April","Mei","Juni","Juli","Agustus","September","Oktober","November","Desember"]

  return (month[monthNumber])
bulan = ""

L = instaloader.Instaloader()
username = "bikecenter.id"

#get post
posts = instaloader.Profile.from_username(L.context, username).get_posts()

#set tanggal
UNTIL = datetime(2022, 12, 1)
SINCE = datetime(2022, 12, 31)

#get post by tanggal
post_datetime = []

for post in takewhile(lambda p: p.date > UNTIL, dropwhile(lambda p: p.date > SINCE, posts)):
    post_datetime.append(post)

#sort by likes
posts_sorted_by_likes = sorted(post_datetime,key=lambda p: p.likes + p.comments, reverse=True)

inout_pptx = r"templates.pptx"
prs = Presentation(inout_pptx)
slide = 4;
indexes = [
  1,2,3
]
count = 0
for index, post in enumerate(posts_sorted_by_likes, 1):

  #download image and caption
  L.download_post(post, target=f"{username}_top_{index}")

  

  #get date
  date = str(post.date)
  l = len(date)
  remove_hour = date[:l-9]

  #get filename
  folder = username+"_top_"+str(index)

  dir_list = os.listdir(folder)
  print(len(dir_list))

  filename = date.replace(" ", "_")
  filename = filename.replace(":", "-")
  filename_txt = filename + "_UTC.txt"
  filename_img = filename + "_UTC.jpg"
  filename_txt = folder+"/"+filename_txt
  filename_img = folder+"/"+filename_img

  tahun, month, tanggal= (int(i) for i in remove_hour.split('-'))

  hari = findDay(remove_hour)
  bulan = calendar.month_name[int(month)]
  like_txt = str(post.likes)
  footer = 'Post ini diupload pada ' + str(hari) + ', ' + str(tanggal) + ' ' + str(bulan) + ' '+ str(tahun) + ' dan mendapat ' + str(like_txt) + ' likes.'

  
  data = [
    "TOP 3 ORGANIC POST BASED ON LIKES",
  ]
  file = open(filename_txt,"r",encoding="utf-8")
  caption = file.read()
  data.append(caption)
  data.append(footer)

  for i,shape,dataisi in zip(indexes,prs.slides[slide].shapes,data):
    if not shape.has_text_frame:
      continue
    shape.text_frame.text = dataisi

    for paragraph in shape.text_frame.paragraphs:
      paragraph.font.color.rgb = RGBColor(255,255,255)
      paragraph.font.size = Pt(10.5)
  

  slide += 1
  file.close()
  count += 1
  if count == 3:
    break

posts_sorted_by_likes_desc = sorted(post_datetime,key=lambda p: p.likes + p.comments, reverse=False)

fileppt = "Monthly Report " + username + " " + bulan + ".pptx"
prs.save(fileppt)



