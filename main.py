import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

inout_pptx = r"templates.pptx"

prs = Presentation(inout_pptx)
slide = prs.slides[0]

top_slide_post_based_like = prs.slides[4]

list_data = [
  "Social Media Report test"
]

for shape,data in zip(slide.shapes, list_data):
  if not shape.has_text_frame:
    continue
  shape.text_frame.text = data



file = open("530jakarta_top_1/2022-12-14_09-56-39_UTC.txt","r",encoding="utf-8")
datatoplikes = [
  "TOP 3 ORGANIC POST BASED ON LIKES"
]
data = file.read()
datatoplikes.append(data)
datatoplikes.append("Post ini diupload pada Rabu, 7 Desember 2022 dan mendapat 44 likes.")

data_into_list = data.split("\n")
indexes = [
  1,2,3
]

for index,shape,data in zip(indexes,top_slide_post_based_like.shapes,datatoplikes):
  if not shape.has_text_frame:
    continue
  shape.text_frame.text = data

  for paragraph in shape.text_frame.paragraphs:
    paragraph.font.color.rgb = RGBColor(255,255,255)
    paragraph.font.size = Pt(10.5)
    # paragraph.text.font.name = 'Roboto'

  # shape.text_frame.paragraphs[0].font.size = Pt(10.5)
  # if index != 1:
  #   shape.text_frame.text.font.name = "Roboto"
  #   shape.text_frame.text.font.size = Pt(10.5)
  # index += 1
# top_slide_post_based_like.shapes.text_frame.text = data_into_list[0]


file.close()
prs.save("test.pptx")