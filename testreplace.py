import collections
import collections.abc
from pptx import Presentation
from pptx_replace import *

prs = Presentation("templates.pptx")
# replace_text(prs, "{Main title}", "this is main report title")
slide = prs.slides[4]
replace_text(slide, slide.shapes[1].text_frame.text, "This is a title test")
# replace_picture(prs.slides[0], "530jakarta_top_2/2022-12-07_10-20-47_UTC.jpg", auto_reshape=True)

print(slide)
prs.save("hasil.pptx")